import csv
import argparse

import pptx
from pptx.util import Pt, Cm


BLANK_SLIDE = 6
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
TEMPLATE_NAME = 'template.pptx'

FONT_NAME = 'Calibri'
FONT_SIZE = Pt(14)


def generate_table(slide, data_rows, column_names):
    # Setup table
    num_rows = len(data_rows)+1
    row_height = Cm(0.78)

    num_cols = 5
    col_widths = list(map(Cm, [3.15, 2.95, 4, 4.15, 3.5]))

    top, left = 476130, 2373459
    tablewidth = sum(col_widths)
    tableheight = (num_rows+1)*row_height

    tableshape = slide.shapes.add_table(
        num_rows, num_cols, left, top, tablewidth, tableheight
    )
    table = tableshape.table

    # Insert table content
    for cell_title, cell in zip(column_names, table.rows[0].cells):
        frame = cell.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = cell_title

        font = run.font
        font.name = FONT_NAME
        font.size = FONT_SIZE
        font.bold = True

    for i, row_data in enumerate(data_rows):
        cellrow = table.rows[i+1]
        for cell_data, cell in zip(row_data.values(), cellrow.cells):
            frame = cell.text_frame
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = cell_data

            font = run.font
            font.name = FONT_NAME
            font.size = FONT_SIZE

    # Format table
    for col_width, column in zip(col_widths, table.columns):
        column.width = col_width


def generate_presentation(pres, data_rows, column_names):
    # Setup slide
    pres.slide_width = SLIDE_WIDTH
    pres.slide_height = SLIDE_HEIGHT

    if len(pres.slides) > 0:
        slide = pres.slides[0]
    else:
        slide = pres.slides.add_slide(pres.slide_layouts[BLANK_SLIDE])

    generate_table(slide, data_rows, column_names)
    return pres


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument("csv_file")
    parser.add_argument("output_file")
    parser.add_argument("-t", "--template", default="template.pptx")

    args = parser.parse_args()
    # Get slide data
    with open(args.csv_file) as f:
        reader = csv.DictReader(f)
        data_rows = [row for row in reader]
        column_names = list(data_rows[0].keys())

    pres = pptx.Presentation(args.template)
    pres = generate_presentation(pres, data_rows, column_names)
    pres.save(args.output_file)
